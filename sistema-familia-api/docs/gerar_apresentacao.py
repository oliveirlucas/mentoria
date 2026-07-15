"""
Gera a apresentacao 'Descomplicando a Arquitetura - Sistema Família API'
seguindo o mesmo estilo visual do template 'Descomplicando APIs.pptx'.

Estilo do template (extraído em _analisar_template.py):
- Slide 16:9 (13.33 x 7.5 in)
- Fontes: Outfit (títulos, bold) / Inter (body)
- Paleta:
    slate-900 #0F172A (título principal)
    slate-800 #1E293B (título pagina)
    slate-600 #475569 (body)
    slate-500 #64748B (subtítulo)
    sky-500   #0EA5E9 (destaque azul)
    sky-400   #38BDF8 (título section)
    emerald-600 #059669 (sucesso)
    amber-600 #D97706 (warning)
    red-600     #DC2626 (erro)
    indigo-600  #4F46E5 (link)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ----------------------------------------------------------------------------
# CORES
# ----------------------------------------------------------------------------
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_800 = RGBColor(0x1E, 0x29, 0x3B)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_600 = RGBColor(0x47, 0x55, 0x69)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
SLATE_200 = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_50 = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SKY_600 = RGBColor(0x02, 0x84, 0xC7)
SKY_500 = RGBColor(0x0E, 0xA5, 0xE9)
SKY_400 = RGBColor(0x38, 0xBD, 0xF8)
SKY_100 = RGBColor(0xE0, 0xF2, 0xFE)
SKY_50 = RGBColor(0xF0, 0xF9, 0xFF)
EMERALD_600 = RGBColor(0x05, 0x96, 0x69)
EMERALD_100 = RGBColor(0xD1, 0xFA, 0xE5)
AMBER_600 = RGBColor(0xD9, 0x77, 0x06)
AMBER_100 = RGBColor(0xFE, 0xF3, 0xC7)
RED_600 = RGBColor(0xDC, 0x26, 0x26)
RED_100 = RGBColor(0xFE, 0xE2, 0xE2)
INDIGO_600 = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_100 = RGBColor(0xE0, 0xE7, 0xFF)
VIOLET_600 = RGBColor(0x7C, 0x3A, 0xED)
VIOLET_100 = RGBColor(0xED, 0xE9, 0xFE)

# ----------------------------------------------------------------------------
# FONTES
# ----------------------------------------------------------------------------
FONT_TITLE = "Outfit"
FONT_BODY = "Inter"
FONT_CODE = "Consolas"

# ----------------------------------------------------------------------------
# SETUP DA APRESENTACAO
# ----------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK_LAYOUT = prs.slide_layouts[6]  # blank


# ----------------------------------------------------------------------------
# HELPERS
# ----------------------------------------------------------------------------
def add_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)


def set_background(slide, color=WHITE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _emu(v):
    """Garante que o valor seja um EMU inteiro (PowerPoint recusa floats)."""
    return int(round(v))


def add_rect(slide, x, y, w, h, fill=WHITE, line=None, line_width=0, shadow=False, rounded=False, corner_radius=0.05):
    if w <= 0 or h <= 0:
        raise ValueError(f"add_rect com dimensao inválida: w={w}, h={h}")
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, _emu(x), _emu(y), _emu(w), _emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    if rounded:
        try:
            shape.adjustments[0] = corner_radius
        except Exception:
            pass
    if not shadow:
        # remove shadow
        sp = shape.shadow
        try:
            spPr = shape._element.spPr
            effectLst = spPr.find(qn("a:effectLst"))
            if effectLst is not None:
                spPr.remove(effectLst)
            effectLst = etree.SubElement(spPr, qn("a:effectLst"))
        except Exception:
            pass
    return shape


def add_line(slide, x1, y1, x2, y2, color=SLATE_200, width=1.0):
    line = slide.shapes.add_connector(1, _emu(x1), _emu(y1), _emu(x2), _emu(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    font=FONT_BODY,
    size=14,
    bold=False,
    color=SLATE_600,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.15,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_rich_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """
    runs: lista de listas. Cada sub-lista representa um parágrafo.
    Cada item da sub-lista é um dict:
      {"text": "...", "font": FONT_BODY, "size": 14, "bold": False, "color": SLATE_600}
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    for pi, paragraph_runs in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for run_spec in paragraph_runs:
            r = p.add_run()
            r.text = run_spec.get("text", "")
            r.font.name = run_spec.get("font", FONT_BODY)
            r.font.size = Pt(run_spec.get("size", 14))
            r.font.bold = run_spec.get("bold", False)
            color = run_spec.get("color", SLATE_600)
            if color is not None:
                r.font.color.rgb = color
    return tb


def add_page_header(slide, section_kicker, title):
    """Adiciona o cabeçalho padrão: subtítulo pequeno azul + título grande."""
    add_text(
        slide,
        Inches(0.6),
        Inches(0.45),
        Inches(12.0),
        Inches(0.35),
        section_kicker.upper(),
        font=FONT_BODY,
        size=12,
        bold=True,
        color=SKY_500,
    )
    add_text(
        slide,
        Inches(0.6),
        Inches(0.80),
        Inches(12.0),
        Inches(0.9),
        title,
        font=FONT_TITLE,
        size=32,
        bold=True,
        color=SLATE_800,
    )
    # linha divisoria fina
    add_line(
        slide,
        Inches(0.6),
        Inches(1.75),
        Inches(12.7),
        Inches(1.75),
        color=SLATE_200,
        width=1.0,
    )


def add_footer(slide, page_num, total):
    """Rodapé com marca e paginação."""
    add_text(
        slide,
        Inches(0.6),
        Inches(7.10),
        Inches(6.0),
        Inches(0.3),
        "Descomplicando a Arquitetura - Sistema Família API",
        font=FONT_BODY,
        size=9,
        color=SLATE_400,
    )
    add_text(
        slide,
        Inches(7.0),
        Inches(7.10),
        Inches(6.0),
        Inches(0.3),
        f"{page_num:02d} / {total:02d}",
        font=FONT_BODY,
        size=9,
        color=SLATE_400,
        align=PP_ALIGN.RIGHT,
    )


def add_code_block(slide, x, y, w, h, lines, size=11):
    """Bloco escuro estilo IDE com linhas de código."""
    add_rect(slide, x, y, w, h, fill=SLATE_900, rounded=True, corner_radius=0.04)
    # tres circulos macOS
    r = Inches(0.10)
    cy = y + Inches(0.18)
    for i, c in enumerate([RGBColor(0xEF, 0x44, 0x44), RGBColor(0xF5, 0x9E, 0x0B), RGBColor(0x22, 0xC5, 0x5E)]):
        cx = x + Inches(0.25) + Inches(0.22) * i
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, r, r)
        dot.fill.solid()
        dot.fill.fore_color.rgb = c
        dot.line.fill.background()

    # linha separadora
    add_line(
        slide,
        x + Inches(0.02),
        y + Inches(0.45),
        x + w - Inches(0.02),
        y + Inches(0.45),
        color=SLATE_700,
        width=0.5,
    )

    # texto
    tb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.55), w - Inches(0.5), h - Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        if isinstance(spec, str):
            r = p.add_run()
            r.text = spec
            r.font.name = FONT_CODE
            r.font.size = Pt(size)
            r.font.color.rgb = SLATE_100
        else:
            # lista de dicts {text, color}
            for run_spec in spec:
                r = p.add_run()
                r.text = run_spec.get("text", "")
                r.font.name = FONT_CODE
                r.font.size = Pt(run_spec.get("size", size))
                r.font.bold = run_spec.get("bold", False)
                r.font.color.rgb = run_spec.get("color", SLATE_100)


def add_pill(slide, x, y, w, h, text, bg=SKY_100, fg=SKY_600, size=11, bold=True, font=FONT_BODY):
    add_rect(slide, x, y, w, h, fill=bg, rounded=True, corner_radius=0.5)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = fg


def add_card(
    slide,
    x,
    y,
    w,
    h,
    title,
    body,
    accent=SKY_500,
    accent_bg=SKY_50,
    title_size=16,
    body_size=12,
    title_color=None,
    show_bar=True,
):
    """Card branco com borda + faixa colorida a esquerda + título + body."""
    # sombra suave via rectangulo cinza atras
    add_rect(slide, x + Inches(0.05), y + Inches(0.06), w, h, fill=SLATE_100, rounded=True, corner_radius=0.06)
    # card em si
    card = add_rect(slide, x, y, w, h, fill=WHITE, line=SLATE_200, line_width=0.75, rounded=True, corner_radius=0.06)
    if show_bar:
        add_rect(slide, x, y, Inches(0.08), h, fill=accent, rounded=False)

    pad_left = Inches(0.30 if show_bar else 0.20)
    pad_top = Inches(0.20)
    add_text(
        slide,
        x + pad_left,
        y + pad_top,
        w - pad_left - Inches(0.15),
        Inches(0.35),
        title,
        font=FONT_TITLE,
        size=title_size,
        bold=True,
        color=title_color if title_color else accent,
    )
    add_text(
        slide,
        x + pad_left,
        y + pad_top + Inches(0.42),
        w - pad_left - Inches(0.15),
        h - pad_top - Inches(0.42) - Inches(0.15),
        body,
        font=FONT_BODY,
        size=body_size,
        color=SLATE_600,
        line_spacing=1.28,
    )


def add_bullet_list(slide, x, y, w, h, items, size=13, color=SLATE_600, bullet_color=None, bullet_char="\u2022", line_spacing=1.35):
    """Cada item pode ser string ou tuple (bold_prefix, resto)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    if bullet_color is None:
        bullet_color = SKY_500

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet
        rb = p.add_run()
        rb.text = f"{bullet_char}  "
        rb.font.name = FONT_BODY
        rb.font.size = Pt(size)
        rb.font.bold = True
        rb.font.color.rgb = bullet_color

        if isinstance(item, tuple):
            head, tail = item
            rh = p.add_run()
            rh.text = head
            rh.font.name = FONT_BODY
            rh.font.size = Pt(size)
            rh.font.bold = True
            rh.font.color.rgb = SLATE_700
            if tail:
                rt = p.add_run()
                rt.text = tail
                rt.font.name = FONT_BODY
                rt.font.size = Pt(size)
                rt.font.color.rgb = color
        else:
            rt = p.add_run()
            rt.text = item
            rt.font.name = FONT_BODY
            rt.font.size = Pt(size)
            rt.font.color.rgb = color


def add_arrow(slide, x1, y1, x2, y2, color=SKY_500, width=2.5):
    conn = slide.shapes.add_connector(1, _emu(x1), _emu(y1), _emu(x2), _emu(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    # setinha
    ln = conn.line
    lnEl = ln._get_or_add_ln()
    tailEnd = etree.SubElement(lnEl, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("h", "med")
    return conn


# ----------------------------------------------------------------------------
# SLIDES
# ----------------------------------------------------------------------------
TOTAL_SLIDES = 27  # atualizar se mudar


def slide_capa():
    slide = add_slide()
    set_background(slide, WHITE)

    # bloco azul decorativo no canto superior direito
    add_rect(slide, Inches(9.5), Inches(-1.0), Inches(6.0), Inches(4.0), fill=SKY_50, rounded=True, corner_radius=0.1)
    add_rect(slide, Inches(10.5), Inches(3.2), Inches(3.5), Inches(3.5), fill=SKY_100, rounded=True, corner_radius=0.1)
    # circulo destaque
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.4), Inches(0.6), Inches(1.6), Inches(1.6))
    circ.fill.solid()
    circ.fill.fore_color.rgb = SKY_500
    circ.line.fill.background()

    # etiqueta
    add_pill(
        slide,
        Inches(0.7),
        Inches(1.15),
        Inches(3.5),
        Inches(0.42),
        "AULA DE ARQUITETURA  |  ALTZEN PRO",
        bg=SKY_100,
        fg=SKY_600,
        size=11,
    )

    # titulo enorme com destaque em azul
    add_rich_text(
        slide,
        Inches(0.6),
        Inches(1.85),
        Inches(11.0),
        Inches(2.4),
        [
            [
                {"text": "Descomplicando ", "font": FONT_TITLE, "size": 56, "bold": True, "color": SLATE_900},
                {"text": "a Arquitetura", "font": FONT_TITLE, "size": 56, "bold": True, "color": SKY_500},
            ],
            [
                {"text": "do Sistema Família API", "font": FONT_TITLE, "size": 44, "bold": True, "color": SLATE_900},
            ],
        ],
        line_spacing=1.05,
    )

    add_text(
        slide,
        Inches(0.7),
        Inches(4.55),
        Inches(11.0),
        Inches(0.6),
        "Do menu de terminal ao Spring Boot: entendendo camadas, MVC, DTOs e injeção de dependência",
        font=FONT_BODY,
        size=17,
        color=SLATE_500,
    )

    # linha divisoria + creditos
    add_line(slide, Inches(0.7), Inches(5.6), Inches(4.5), Inches(5.6), color=SKY_500, width=2.0)
    add_text(
        slide,
        Inches(0.7),
        Inches(5.75),
        Inches(8.0),
        Inches(0.4),
        "Aula preparada para o aluno na Mentoria AltzenPro",
        font=FONT_BODY,
        size=13,
        bold=True,
        color=SLATE_700,
    )
    add_text(
        slide,
        Inches(0.7),
        Inches(6.10),
        Inches(8.0),
        Inches(0.4),
        "Baseada no projeto sistema-família-api  -  Spring Boot 3 + Java 17",
        font=FONT_BODY,
        size=12,
        color=SLATE_500,
    )


def slide_section(kicker, title, subtitle):
    slide = add_slide()
    set_background(slide, SLATE_50)

    # bloco decorativo lateral
    add_rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, fill=SKY_500, rounded=False)

    add_pill(
        slide,
        Inches(0.9),
        Inches(2.3),
        Inches(2.2),
        Inches(0.45),
        kicker,
        bg=SKY_100,
        fg=SKY_600,
        size=12,
    )
    add_text(
        slide,
        Inches(0.9),
        Inches(2.95),
        Inches(12.0),
        Inches(1.4),
        title,
        font=FONT_TITLE,
        size=54,
        bold=True,
        color=SLATE_900,
    )
    add_text(
        slide,
        Inches(0.9),
        Inches(4.30),
        Inches(11.5),
        Inches(0.7),
        subtitle,
        font=FONT_BODY,
        size=18,
        color=SLATE_500,
    )
    add_line(slide, Inches(0.9), Inches(5.05), Inches(3.9), Inches(5.05), color=SKY_500, width=2.5)


def slide_de_onde_viemos():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Ponto de partida", "De onde viemos: do console para a API")

    add_text(
        slide,
        Inches(0.6),
        Inches(1.95),
        Inches(12.0),
        Inches(0.5),
        "O projeto original rodava no terminal, com um menu e um loop while. Vamos ver o que muda:",
        font=FONT_BODY,
        size=14,
        color=SLATE_600,
    )

    # Duas colunas: Console vs API
    col_w = Inches(6.0)
    left_x = Inches(0.6)
    right_x = Inches(6.75)
    top_y = Inches(2.65)
    col_h = Inches(4.15)

    # Console
    add_rect(slide, left_x, top_y, col_w, col_h, fill=SLATE_50, rounded=True, corner_radius=0.04, line=SLATE_200, line_width=0.75)
    add_pill(slide, left_x + Inches(0.3), top_y + Inches(0.25), Inches(1.8), Inches(0.4), "PROJETO CONSOLE", bg=SLATE_200, fg=SLATE_700, size=10)
    add_text(
        slide,
        left_x + Inches(0.3),
        top_y + Inches(0.75),
        col_w - Inches(0.6),
        Inches(0.5),
        "sistema-família (Java puro)",
        font=FONT_TITLE, size=20, bold=True, color=SLATE_800,
    )
    add_bullet_list(
        slide,
        left_x + Inches(0.3),
        top_y + Inches(1.35),
        col_w - Inches(0.6),
        Inches(2.8),
        [
            ("Main.java ", "com while(continuar) e if/else por opção"),
            ("ArrayList ", "declarados dentro do main() (dados locais)"),
            ("GestorPatrimonio ", "com métodos static (sem instância)"),
            ("Scanner + System.out ", "para conversar com o usuário"),
            ("if (nome.isBlank()) ... ", "validações espalhadas no código"),
        ],
        size=12,
        bullet_color=SLATE_500,
    )

    # API
    add_rect(slide, right_x, top_y, col_w, col_h, fill=SKY_50, rounded=True, corner_radius=0.04, line=SKY_100, line_width=0.75)
    add_pill(slide, right_x + Inches(0.3), top_y + Inches(0.25), Inches(1.6), Inches(0.4), "PROJETO API", bg=SKY_500, fg=WHITE, size=10)
    add_text(
        slide,
        right_x + Inches(0.3),
        top_y + Inches(0.75),
        col_w - Inches(0.6),
        Inches(0.5),
        "sistema-família-api (Spring Boot)",
        font=FONT_TITLE, size=20, bold=True, color=SKY_600,
    )
    add_bullet_list(
        slide,
        right_x + Inches(0.3),
        top_y + Inches(1.35),
        col_w - Inches(0.6),
        Inches(2.8),
        [
            ("Controllers ", "REST: um endpoint para cada ação"),
            ("Repositories ", "singleton com lista em memória"),
            ("Services ", "com a mesma lógica de negócio"),
            ("HTTP + JSON ", "no lugar de Scanner e System.out"),
            ("@NotBlank / @Positive ", "validação declarativa via Jakarta"),
        ],
        size=12,
        bullet_color=SKY_500,
    )

    add_footer(slide, 3, TOTAL_SLIDES)


def slide_o_que_e_mvc():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Fundamentos", "O que é MVC (e por que ele importa)?")

    add_text(
        slide,
        Inches(0.6),
        Inches(1.95),
        Inches(12.0),
        Inches(0.5),
        "MVC = Model-View-Controller. É um padrão que divide responsabilidades para o código não virar bagunça.",
        font=FONT_BODY,
        size=14,
        color=SLATE_600,
    )

    # Tres cards horizontais: M, V, C
    y = Inches(2.75)
    h = Inches(3.4)
    gap = Inches(0.25)
    w = (SLIDE_W - Inches(1.2) - gap * 2) / 3
    x1 = Inches(0.6)
    x2 = x1 + w + gap
    x3 = x2 + w + gap

    add_card(
        slide, x1, y, w, h,
        title="M - Model",
        body=(
            "Representa os DADOS e as regras básicas do domínio.\n\n"
            "Aqui: Pessoa, Animal, Casa, Carro.\n\n"
            "Não sabe nada sobre HTTP, banco ou tela."
        ),
        accent=INDIGO_600,
        accent_bg=INDIGO_100,
        title_color=INDIGO_600,
    )
    add_card(
        slide, x2, y, w, h,
        title="V - View",
        body=(
            "É o que o usuário vê.\n\n"
            "Numa API REST a View é o próprio JSON de resposta (Response DTOs).\n\n"
            "Em uma app web seria HTML/CSS/JS."
        ),
        accent=VIOLET_600,
        accent_bg=VIOLET_100,
        title_color=VIOLET_600,
    )
    add_card(
        slide, x3, y, w, h,
        title="C - Controller",
        body=(
            "Recebe o pedido, coordena o que fazer e devolve a resposta.\n\n"
            "Não contém regras de negócio: ele DELEGA para o Service.\n\n"
            "É a porta de entrada da API."
        ),
        accent=SKY_500,
        accent_bg=SKY_100,
        title_color=SKY_500,
    )

    add_text(
        slide,
        Inches(0.6),
        Inches(6.55),
        Inches(12.0),
        Inches(0.4),
        "Spring adiciona duas camadas extras (Service e Repository) que deixam o código ainda mais organizado.",
        font=FONT_BODY, size=12, color=SLATE_500,
    )

    add_footer(slide, 4, TOTAL_SLIDES)


def slide_arquitetura_camadas():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Arquitetura", "Como as camadas conversam entre si")

    add_text(
        slide,
        Inches(0.6),
        Inches(1.95),
        Inches(12.0),
        Inches(0.5),
        "Cada requisição HTTP atravessa uma pilha de camadas. Cada uma tem UMA responsabilidade.",
        font=FONT_BODY, size=14, color=SLATE_600,
    )

    # 5 caixas horizontais + arrows
    box_w = Inches(2.15)
    box_h = Inches(1.2)
    y = Inches(3.15)
    gap = Inches(0.15)
    total_w = box_w * 5 + gap * 4
    start_x = (SLIDE_W - total_w) / 2

    boxes = [
        ("Cliente", "Browser, Postman, app...", SLATE_700, SLATE_100),
        ("Controller", "@RestController\nRecebe HTTP / JSON", SKY_500, SKY_50),
        ("Service", "@Service\nRegras de negócio", INDIGO_600, INDIGO_100),
        ("Repository", "@Repository\nAcesso aos dados", VIOLET_600, VIOLET_100),
        ("Modelo em memória", "List<Pessoa> singleton", SLATE_700, SLATE_100),
    ]

    positions = []
    for i, (title, body, accent, bg) in enumerate(boxes):
        x = start_x + (box_w + gap) * i
        add_rect(slide, x, y, box_w, box_h, fill=bg, rounded=True, corner_radius=0.10, line=accent, line_width=1.25)
        add_text(slide, x + Inches(0.1), y + Inches(0.15), box_w - Inches(0.2), Inches(0.35),
                 title, font=FONT_TITLE, size=15, bold=True, color=accent, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.1), y + Inches(0.55), box_w - Inches(0.2), Inches(0.6),
                 body, font=FONT_BODY, size=10, color=SLATE_600, align=PP_ALIGN.CENTER, line_spacing=1.15)
        positions.append((x, y))

    # setas
    for i in range(4):
        x1 = positions[i][0] + box_w
        x2 = positions[i + 1][0]
        y_arrow = y + box_h / 2
        add_arrow(slide, x1 - Inches(0.02), y_arrow, x2 + Inches(0.02), y_arrow, color=SLATE_400, width=1.8)

    # Legenda abaixo
    add_text(
        slide,
        Inches(0.6),
        Inches(4.85),
        Inches(12.0),
        Inches(0.35),
        "E a volta? A resposta faz o caminho inverso, convertida em DTO → JSON.",
        font=FONT_BODY, size=12, color=SLATE_500, align=PP_ALIGN.CENTER,
    )

    # Cards de DTO e Exception
    y2 = Inches(5.55)
    h2 = Inches(1.35)
    w_mini = Inches(5.9)
    add_card(
        slide, Inches(0.6), y2, w_mini, h2,
        title="DTO (Data Transfer Object)",
        body=(
            "Objeto que ENTRA (Request) ou SAI (Response) da API. Evita expor o Model interno "
            "e garante um contrato claro no JSON."
        ),
        accent=EMERALD_600, title_color=EMERALD_600,
        title_size=15, body_size=11,
    )
    add_card(
        slide, Inches(6.85), y2, w_mini, h2,
        title="Exception + Handler",
        body=(
            "Erros viram exceções (BusinessException, NotFoundException). O GlobalExceptionHandler "
            "traduz cada uma em um JSON com status HTTP correto."
        ),
        accent=RED_600, title_color=RED_600,
        title_size=15, body_size=11,
    )

    add_footer(slide, 5, TOTAL_SLIDES)


def slide_camadas_pratica_header():
    slide_section(
        "PARTE 2",
        "As Camadas na Prática",
        "Vamos abrir cada uma: Model, Repository, Service, Controller, DTO e Exception.",
    )


def slide_model():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Camada 01", "Model - as classes que representam o domínio")

    # coluna esquerda: explicacao
    add_bullet_list(
        slide,
        Inches(0.6),
        Inches(2.1),
        Inches(5.9),
        Inches(4.5),
        [
            ("O que é? ", "classes que representam CONCEITOS do mundo real (Pessoa, Casa, Carro)."),
            ("São POJOs: ", "somente atributos + getters/setters + construtor. Não usam Spring."),
            ("Encapsulamento: ", "atributos private, acesso via métodos. Isso protege os dados."),
            ("Id: ", "cada entidade ganhou um campo id para poder ser identificada nas rotas."),
            ("Não conhece: ", "HTTP, JSON, banco de dados, tela. É puro domínio."),
        ],
        size=13,
    )

    # coluna direita: codigo
    add_code_block(
        slide,
        Inches(6.75),
        Inches(2.1),
        Inches(6.05),
        Inches(4.5),
        [
            [{"text": "public class ", "color": SKY_400},
             {"text": "Pessoa", "color": WHITE, "bold": True},
             {"text": " {", "color": SLATE_100}],
            "",
            [{"text": "    private Long ", "color": SKY_400},
             {"text": "id;", "color": SLATE_100}],
            [{"text": "    private String ", "color": SKY_400},
             {"text": "nome;", "color": SLATE_100}],
            [{"text": "    private int ", "color": SKY_400},
             {"text": "idade;", "color": SLATE_100}],
            [{"text": "    private String ", "color": SKY_400},
             {"text": "parentesco;", "color": SLATE_100}],
            [{"text": "    private double ", "color": SKY_400},
             {"text": "carteira;", "color": SLATE_100}],
            [{"text": "    private double ", "color": SKY_400},
             {"text": "dívida;", "color": SLATE_100}],
            "",
            [{"text": "    public ", "color": SKY_400},
             {"text": "Pessoa", "color": WHITE, "bold": True},
             {"text": "(String nome, int idade, ...) { ... }", "color": SLATE_100}],
            "",
            [{"text": "    // getters e setters", "color": SLATE_400}],
            "}",
        ],
        size=11,
    )

    add_footer(slide, 7, TOTAL_SLIDES)


def slide_repository():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Camada 02", "Repository - onde os dados ficam guardados")

    add_bullet_list(
        slide,
        Inches(0.6),
        Inches(2.1),
        Inches(5.9),
        Inches(4.5),
        [
            ("Anotação: ", "@Repository (o Spring identifica como componente injetável)."),
            ("Papel: ", "isolar a forma de armazenar dados do resto do sistema."),
            ("Aqui é simples: ", "uma List em memória + gerador de id (AtomicLong)."),
            ("Singleton: ", "o Spring cria UMA instância só, compartilhada por toda a aplicação."),
            ("Amanhã... ", "trocamos por Spring Data JPA + banco e o Service nem percebe."),
        ],
        size=13,
    )

    add_code_block(
        slide,
        Inches(6.75),
        Inches(2.1),
        Inches(6.05),
        Inches(4.5),
        [
            [{"text": "@Repository", "color": EMERALD_100}],
            [{"text": "public class ", "color": SKY_400},
             {"text": "PessoaRepository", "color": WHITE, "bold": True},
             {"text": " {", "color": SLATE_100}],
            "",
            [{"text": "    private final List<Pessoa> ", "color": SLATE_100},
             {"text": "pessoas ", "color": WHITE},
             {"text": "= ", "color": SLATE_100},
             {"text": "new ", "color": SKY_400},
             {"text": "ArrayList<>();", "color": SLATE_100}],
            [{"text": "    private final AtomicLong ", "color": SLATE_100},
             {"text": "seq ", "color": WHITE},
             {"text": "= ", "color": SLATE_100},
             {"text": "new ", "color": SKY_400},
             {"text": "AtomicLong(1);", "color": SLATE_100}],
            "",
            [{"text": "    public ", "color": SKY_400},
             {"text": "Pessoa ", "color": SLATE_100},
             {"text": "salvar", "color": WHITE, "bold": True},
             {"text": "(Pessoa p) {", "color": SLATE_100}],
            [{"text": "        p.setId(seq.getAndIncrement());", "color": SLATE_100}],
            [{"text": "        pessoas.add(p);", "color": SLATE_100}],
            [{"text": "        return ", "color": SKY_400},
             {"text": "p;", "color": SLATE_100}],
            [{"text": "    }", "color": SLATE_100}],
            "",
            [{"text": "    // buscarPorId, listarTodos, ...", "color": SLATE_400}],
            "}",
        ],
        size=11,
    )

    add_footer(slide, 8, TOTAL_SLIDES)


def slide_service():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Camada 03", "Service - o cérebro das regras de negócio")

    add_bullet_list(
        slide,
        Inches(0.6),
        Inches(2.1),
        Inches(5.9),
        Inches(4.5),
        [
            ("Anotação: ", "@Service (semântica de 'regra de negócio')."),
            ("O que faz: ", "orquestra Repositories e valida regras que envolvem o domínio."),
            ("Não conhece HTTP: ", "nada de HttpRequest, JSON ou status code aqui."),
            ("Substitui o antigo: ", "GestorPatrimonio (que era static) virou GestorPatrimonioService."),
            ("Quando algo da errado: ", "lança BusinessException com uma mensagem clara."),
        ],
        size=13,
    )

    add_code_block(
        slide,
        Inches(6.75),
        Inches(2.1),
        Inches(6.05),
        Inches(4.5),
        [
            [{"text": "@Service", "color": EMERALD_100}],
            [{"text": "public class ", "color": SKY_400},
             {"text": "PessoaService", "color": WHITE, "bold": True},
             {"text": " {", "color": SLATE_100}],
            "",
            [{"text": "    private final PessoaRepository ", "color": SLATE_100},
             {"text": "repository;", "color": WHITE}],
            "",
            [{"text": "    public ", "color": SKY_400},
             {"text": "PessoaService", "color": WHITE, "bold": True},
             {"text": "(PessoaRepository r) {", "color": SLATE_100}],
            [{"text": "        this.repository = r;   ", "color": SLATE_100},
             {"text": "// injeção via construtor", "color": SLATE_400}],
            [{"text": "    }", "color": SLATE_100}],
            "",
            [{"text": "    public ", "color": SKY_400},
             {"text": "Pessoa ", "color": SLATE_100},
             {"text": "cadastrar", "color": WHITE, "bold": True},
             {"text": "(Pessoa p) {", "color": SLATE_100}],
            [{"text": "        if ", "color": SKY_400},
             {"text": "(p.getIdade() < 0)", "color": SLATE_100}],
            [{"text": "            throw new ", "color": SKY_400},
             {"text": "BusinessException(", "color": SLATE_100},
             {"text": "\"Idade inválida\"", "color": AMBER_100},
             {"text": ");", "color": SLATE_100}],
            [{"text": "        return ", "color": SKY_400},
             {"text": "repository.salvar(p);", "color": SLATE_100}],
            [{"text": "    }", "color": SLATE_100}],
            "}",
        ],
        size=10,
    )

    add_footer(slide, 9, TOTAL_SLIDES)


def slide_controller():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Camada 04", "Controller - a porta de entrada HTTP")

    add_bullet_list(
        slide,
        Inches(0.6),
        Inches(2.1),
        Inches(5.9),
        Inches(4.5),
        [
            ("Anotação: ", "@RestController + @RequestMapping('/pessoas')."),
            ("Roteia: ", "@GetMapping, @PostMapping, @PutMapping, @DeleteMapping."),
            ("Recebe: ", "@RequestBody (JSON -> DTO) e @PathVariable (id na URL)."),
            ("Valida: ", "@Valid dispara as anotações do Jakarta antes de entrar no método."),
            ("Devolve: ", "ResponseEntity com código HTTP + JSON de resposta."),
            ("Não contém regra de negócio: ", "APENAS delega para o Service."),
        ],
        size=13,
    )

    add_code_block(
        slide,
        Inches(6.75),
        Inches(2.1),
        Inches(6.05),
        Inches(4.5),
        [
            [{"text": "@RestController", "color": EMERALD_100}],
            [{"text": "@RequestMapping(", "color": EMERALD_100},
             {"text": "\"/pessoas\"", "color": AMBER_100},
             {"text": ")", "color": EMERALD_100}],
            [{"text": "public class ", "color": SKY_400},
             {"text": "PessoaController", "color": WHITE, "bold": True},
             {"text": " {", "color": SLATE_100}],
            "",
            [{"text": "    private final PessoaService ", "color": SLATE_100},
             {"text": "service;", "color": WHITE}],
            "",
            [{"text": "    public ", "color": SKY_400},
             {"text": "PessoaController", "color": WHITE, "bold": True},
             {"text": "(PessoaService s) { this.service = s; }", "color": SLATE_100}],
            "",
            [{"text": "    @PostMapping", "color": EMERALD_100}],
            [{"text": "    public ", "color": SKY_400},
             {"text": "ResponseEntity<PessoaResponse> ", "color": SLATE_100},
             {"text": "cadastrar", "color": WHITE, "bold": True},
             {"text": "(", "color": SLATE_100}],
            [{"text": "            @Valid @RequestBody ", "color": EMERALD_100},
             {"text": "PessoaRequest req) {", "color": SLATE_100}],
            [{"text": "        var salva = service.cadastrar(req.toModel());", "color": SLATE_100}],
            [{"text": "        return ", "color": SKY_400},
             {"text": "ResponseEntity.status(201).body(PessoaResponse.from(salva));", "color": SLATE_100}],
            [{"text": "    }", "color": SLATE_100}],
            "}",
        ],
        size=10,
    )

    add_footer(slide, 10, TOTAL_SLIDES)


def slide_dto():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Camada 05", "DTOs - o contrato de entrada e saída")

    add_text(
        slide,
        Inches(0.6),
        Inches(1.95),
        Inches(12.0),
        Inches(0.5),
        "Por que não mandar o Model direto no JSON? Porque acopla API + banco e vaza dados internos.",
        font=FONT_BODY, size=14, color=SLATE_600,
    )

    # dois blocos lado a lado
    add_card(
        slide, Inches(0.6), Inches(2.65), Inches(6.0), Inches(1.55),
        title="Request (o que ENTRA)",
        body=(
            "Cada @PostMapping / @PutMapping recebe um *Request.\n"
            "É ali que colocamos as validações @NotBlank, @Positive, @Min, etc."
        ),
        accent=SKY_500, title_color=SKY_500,
    )
    add_card(
        slide, Inches(6.85), Inches(2.65), Inches(6.0), Inches(1.55),
        title="Response (o que SAI)",
        body=(
            "Devolvemos SOMENTE os campos que fazem sentido para o cliente.\n"
            "Ex: em vez de expor o objeto 'Pessoa dono', devolvemos 'donoId' e 'donoNome'."
        ),
        accent=EMERALD_600, title_color=EMERALD_600,
    )

    # dois blocos de codigo
    add_code_block(
        slide, Inches(0.6), Inches(4.35), Inches(6.0), Inches(2.5),
        [
            [{"text": "public record ", "color": SKY_400},
             {"text": "PessoaRequest", "color": WHITE, "bold": True},
             {"text": "(", "color": SLATE_100}],
            [{"text": "    @NotBlank ", "color": EMERALD_100},
             {"text": "String nome,", "color": SLATE_100}],
            [{"text": "    @Min(0) ", "color": EMERALD_100},
             {"text": "int idade,", "color": SLATE_100}],
            [{"text": "    @NotBlank ", "color": EMERALD_100},
             {"text": "String parentesco,", "color": SLATE_100}],
            [{"text": "    @PositiveOrZero ", "color": EMERALD_100},
             {"text": "double carteira", "color": SLATE_100}],
            [{"text": ") {}", "color": SLATE_100}],
        ],
        size=11,
    )
    add_code_block(
        slide, Inches(6.85), Inches(4.35), Inches(6.0), Inches(2.5),
        [
            [{"text": "public record ", "color": SKY_400},
             {"text": "PessoaResponse", "color": WHITE, "bold": True},
             {"text": "(", "color": SLATE_100}],
            [{"text": "    Long id,", "color": SLATE_100}],
            [{"text": "    String nome, int idade,", "color": SLATE_100}],
            [{"text": "    String parentesco,", "color": SLATE_100}],
            [{"text": "    double carteira, double dívida", "color": SLATE_100}],
            [{"text": ") {", "color": SLATE_100}],
            [{"text": "    public static ", "color": SKY_400},
             {"text": "PessoaResponse ", "color": SLATE_100},
             {"text": "from", "color": WHITE, "bold": True},
             {"text": "(Pessoa p) { ... }", "color": SLATE_100}],
            [{"text": "}", "color": SLATE_100}],
        ],
        size=11,
    )

    add_footer(slide, 11, TOTAL_SLIDES)


def slide_exception():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Camada 06", "Exception + GlobalExceptionHandler")

    add_bullet_list(
        slide,
        Inches(0.6),
        Inches(2.1),
        Inches(5.9),
        Inches(4.5),
        [
            ("BusinessException: ", "regra de negócio violada (422)."),
            ("NotFoundException: ", "id não encontrado (404)."),
            ("MethodArgumentNotValidException: ", "campo inválido no @Valid (400)."),
            ("GlobalExceptionHandler: ", "anotado com @RestControllerAdvice, intercepta TUDO."),
            ("Vantagem: ", "código do Service permanece limpo, focado na regra."),
            ("Resposta padrão: ", "sempre um JSON no mesmo formato (timestamp, status, erro, mensagem)."),
        ],
        size=13,
    )

    add_code_block(
        slide,
        Inches(6.75),
        Inches(2.1),
        Inches(6.05),
        Inches(4.5),
        [
            [{"text": "@RestControllerAdvice", "color": EMERALD_100}],
            [{"text": "public class ", "color": SKY_400},
             {"text": "GlobalExceptionHandler", "color": WHITE, "bold": True},
             {"text": " {", "color": SLATE_100}],
            "",
            [{"text": "    @ExceptionHandler(BusinessException.class)", "color": EMERALD_100}],
            [{"text": "    public ", "color": SKY_400},
             {"text": "ResponseEntity<ErroResponse> ", "color": SLATE_100},
             {"text": "regra", "color": WHITE, "bold": True},
             {"text": "(BusinessException e) {", "color": SLATE_100}],
            [{"text": "        return ", "color": SKY_400},
             {"text": "ResponseEntity", "color": SLATE_100}],
            [{"text": "            .status(422)", "color": SLATE_100}],
            [{"text": "            .body(ErroResponse.of(422, ", "color": SLATE_100},
             {"text": "\"Regra violada\"", "color": AMBER_100},
             {"text": ", e.getMessage()));", "color": SLATE_100}],
            [{"text": "    }", "color": SLATE_100}],
            "",
            [{"text": "    // handlers para 400, 404, 500...", "color": SLATE_400}],
            "}",
        ],
        size=10,
    )

    add_footer(slide, 12, TOTAL_SLIDES)


def slide_di_header():
    slide_section(
        "PARTE 3",
        "Injeção de Dependência",
        "Como as camadas se encontram sem precisar dar 'new' em nada.",
    )


def slide_di_o_que_e():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Injeção de dependência", "O que é DI e por que ela facilita a vida?")

    # Dois cards: SEM DI vs COM DI
    add_card(
        slide, Inches(0.6), Inches(2.1), Inches(6.0), Inches(4.7),
        title="SEM injeção de dependência",
        body=(
            "Cada classe cria as suas dependências com 'new'.\n\n"
            "* Alto acoplamento: se o construtor mudar, quebra tudo.\n"
            "* Difícil de testar: não da para trocar por um mock.\n"
            "* Duplicação: cada canto do código cria a mesma instância."
        ),
        accent=RED_600, title_color=RED_600,
        title_size=17, body_size=12,
    )
    add_card(
        slide, Inches(6.85), Inches(2.1), Inches(6.0), Inches(4.7),
        title="COM injeção de dependência",
        body=(
            "A classe apenas DECLARA o que precisa. O Spring entrega pronto.\n\n"
            "* Baixo acoplamento: falamos com a interface, não com o 'new'.\n"
            "* Fácil de testar: injeta-se um mock/fake no lugar do real.\n"
            "* Instância única reutilizada (singleton) - performático e coerente."
        ),
        accent=EMERALD_600, title_color=EMERALD_600,
        title_size=17, body_size=12,
    )

    add_footer(slide, 14, TOTAL_SLIDES)


def slide_di_como_spring_faz():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Injeção de dependência", "Como o Spring faz a mágica acontecer")

    # tres passos horizontais
    boxes = [
        ("1. Marca",
         "Colocamos @Component / @Service / @Repository / @RestController nas classes.",
         SKY_500, SKY_50),
        ("2. Escaneia",
         "No startup, o Spring varre os pacotes e registra tudo no ApplicationContext (o 'sacola de beans').",
         INDIGO_600, INDIGO_100),
        ("3. Injeta",
         "Ao criar uma classe, ele olha o construtor e entrega as dependências já prontas.",
         EMERALD_600, EMERALD_100),
    ]
    y = Inches(2.15)
    w = Inches(4.05)
    h = Inches(2.35)
    gap = Inches(0.15)
    total_w = w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2

    for i, (title, body, accent, bg) in enumerate(boxes):
        x = start_x + (w + gap) * i
        add_rect(slide, x, y, w, h, fill=bg, rounded=True, corner_radius=0.06, line=accent, line_width=0.75)
        add_text(slide, x + Inches(0.3), y + Inches(0.25), w - Inches(0.6), Inches(0.5),
                 title, font=FONT_TITLE, size=20, bold=True, color=accent)
        add_text(slide, x + Inches(0.3), y + Inches(0.85), w - Inches(0.6), h - Inches(1.0),
                 body, font=FONT_BODY, size=13, color=SLATE_600, line_spacing=1.35)

    # exemplo de codigo abaixo
    add_text(
        slide,
        Inches(0.6),
        Inches(4.75),
        Inches(12.0),
        Inches(0.35),
        "A forma recomendada hoje é injeção via construtor (final + sem @Autowired explícito):",
        font=FONT_BODY, size=13, color=SLATE_600, bold=True,
    )

    add_code_block(
        slide,
        Inches(0.6),
        Inches(5.20),
        Inches(12.15),
        Inches(1.75),
        [
            [{"text": "@RestController", "color": EMERALD_100}],
            [{"text": "public class ", "color": SKY_400},
             {"text": "PessoaController", "color": WHITE, "bold": True},
             {"text": " {", "color": SLATE_100}],
            [{"text": "    private final PessoaService service;   ", "color": SLATE_100},
             {"text": "// campo 'final' garante que será injetado uma única vez", "color": SLATE_400}],
            [{"text": "    public ", "color": SKY_400},
             {"text": "PessoaController", "color": WHITE, "bold": True},
             {"text": "(PessoaService service) { this.service = service; }", "color": SLATE_100}],
            [{"text": "}", "color": SLATE_100}],
        ],
        size=11,
    )

    add_footer(slide, 15, TOTAL_SLIDES)


def slide_di_no_projeto():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Injeção de dependência", "A cadeia de injeção no nosso projeto")

    # cadeia visual (4 caixas verticais + setas)
    boxes = [
        ("Spring ApplicationContext", "cria e gerencia todos os beans", SLATE_700, SLATE_100),
        ("PessoaRepository", "@Repository - lista em memória", VIOLET_600, VIOLET_100),
        ("PessoaService", "@Service - recebe PessoaRepository", INDIGO_600, INDIGO_100),
        ("PessoaController", "@RestController - recebe PessoaService", SKY_500, SKY_100),
    ]

    y = Inches(2.15)
    w = Inches(4.8)
    h = Inches(1.05)
    gap = Inches(0.30)
    x = (SLIDE_W - w) / 2

    positions = []
    for i, (title, body, accent, bg) in enumerate(boxes):
        add_rect(slide, x, y, w, h, fill=bg, rounded=True, corner_radius=0.10, line=accent, line_width=1.0)
        add_text(slide, x + Inches(0.3), y + Inches(0.15), w - Inches(0.6), Inches(0.4),
                 title, font=FONT_TITLE, size=17, bold=True, color=accent)
        add_text(slide, x + Inches(0.3), y + Inches(0.60), w - Inches(0.6), Inches(0.4),
                 body, font=FONT_BODY, size=12, color=SLATE_600)
        positions.append((x, y))
        y = y + h + gap

    # setas verticais entre elas
    for i in range(3):
        x_center = positions[i][0] + w / 2
        y_start = positions[i][1] + h + Inches(0.02)
        y_end = positions[i + 1][1] - Inches(0.02)
        add_arrow(slide, x_center, y_start, x_center, y_end, color=SLATE_400, width=1.8)

    # legenda a direita
    add_text(
        slide,
        Inches(9.6),
        Inches(2.3),
        Inches(3.5),
        Inches(4.5),
        "Você nunca escreve 'new'.\n"
        "\n"
        "O Spring:\n"
        "1) Cria o Repository.\n"
        "2) Cria o Service e passa o Repository.\n"
        "3) Cria o Controller e passa o Service.\n"
        "\n"
        "Se amanhã algum construtor mudar, o Spring resolve.",
        font=FONT_BODY, size=12, color=SLATE_600, line_spacing=1.35,
    )

    add_footer(slide, 16, TOTAL_SLIDES)


def slide_ciclo_header():
    slide_section(
        "PARTE 4",
        "O Ciclo de uma Requisição",
        "Do clique no Postman até a resposta em JSON, passo a passo.",
    )


def slide_ciclo_passo_a_passo():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Ciclo da requisição", "Passo a passo: POST /pessoas")

    steps = [
        ("1. Cliente envia", "POST /pessoas com JSON no body.", SKY_500),
        ("2. Spring recebe", "Dispatcher acha o Controller certo e o método pelo verbo/URL.", INDIGO_600),
        ("3. Body vira DTO", "Jackson converte JSON em PessoaRequest.", VIOLET_600),
        ("4. @Valid roda", "Validações disparam. Se falha → 400 automático.", AMBER_600),
        ("5. Controller delega", "Converte DTO em Pessoa e chama service.cadastrar(...).", SKY_500),
        ("6. Service aplica regra", "Valida idade/dono/saldo. Chama o Repository.", INDIGO_600),
        ("7. Repository salva", "Gera id e adiciona na lista em memória.", VIOLET_600),
        ("8. Volta a resposta", "Service devolve Pessoa -> Controller monta PessoaResponse -> JSON 201.", EMERALD_600),
    ]

    y = Inches(2.0)
    row_h = Inches(0.55)
    gap = Inches(0.05)

    for i, (title, body, color) in enumerate(steps):
        yy = y + (row_h + gap) * i
        # bolinha numerada
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), yy + Inches(0.02), Inches(0.48), Inches(0.48))
        circ.fill.solid(); circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        add_text(slide, Inches(0.6), yy + Inches(0.02), Inches(0.48), Inches(0.48),
                 str(i + 1), font=FONT_TITLE, size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # texto
        add_text(slide, Inches(1.22), yy + Inches(0.0), Inches(11.5), Inches(0.32),
                 title, font=FONT_TITLE, size=14, bold=True, color=SLATE_800)
        add_text(slide, Inches(1.22), yy + Inches(0.28), Inches(11.5), Inches(0.3),
                 body, font=FONT_BODY, size=11, color=SLATE_600)

    add_footer(slide, 18, TOTAL_SLIDES)


def slide_validacao():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Ciclo da requisição", "Validação declarativa com Jakarta Bean Validation")

    add_bullet_list(
        slide,
        Inches(0.6),
        Inches(2.1),
        Inches(5.9),
        Inches(4.5),
        [
            ("Antes: ", "if (nome == null || nome.isBlank()) { ... } espalhado no código."),
            ("Agora: ", "anotamos o campo do DTO e o Spring valida sozinho."),
            ("@NotBlank ", "- string obrigatória e não vazia."),
            ("@Min(0) / @Max ", "- limites numéricos."),
            ("@Positive / @PositiveOrZero ", "- exige número > 0 (ou >= 0)."),
            ("@Email, @Pattern, @Size ", "- outras validações prontas."),
            ("Se falhar: ", "400 Bad Request automático com a lista de erros."),
        ],
        size=13,
    )

    add_code_block(
        slide,
        Inches(6.75),
        Inches(2.1),
        Inches(6.05),
        Inches(4.5),
        [
            [{"text": "public record ", "color": SKY_400},
             {"text": "CasaRequest", "color": WHITE, "bold": True},
             {"text": "(", "color": SLATE_100}],
            [{"text": "    @NotBlank ", "color": EMERALD_100},
             {"text": "String nome,", "color": SLATE_100}],
            [{"text": "    @Min(0) ", "color": EMERALD_100},
             {"text": "int idade,", "color": SLATE_100}],
            [{"text": "    @NotBlank ", "color": EMERALD_100},
             {"text": "String endereco,", "color": SLATE_100}],
            [{"text": "    @Positive ", "color": EMERALD_100},
             {"text": "double valor,", "color": SLATE_100}],
            [{"text": "    @NotNull ", "color": EMERALD_100},
             {"text": "Long donoId", "color": SLATE_100}],
            [{"text": ") {}", "color": SLATE_100}],
            "",
            [{"text": "// no Controller:", "color": SLATE_400}],
            [{"text": "public ", "color": SKY_400},
             {"text": "ResponseEntity<...> cadastrar(", "color": SLATE_100}],
            [{"text": "    @Valid @RequestBody ", "color": EMERALD_100},
             {"text": "CasaRequest req) { ... }", "color": SLATE_100}],
        ],
        size=11,
    )

    add_footer(slide, 19, TOTAL_SLIDES)


def slide_estrutura_header():
    slide_section(
        "PARTE 5",
        "Estrutura de Pastas e Rotas",
        "Como o código está organizado no repositório e quais URLs a API expõe.",
    )


def slide_estrutura_pastas():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Estrutura de pastas", "Como o projeto está organizado")

    lines = [
        "sistema-família-api/",
        "|-- pom.xml                                  <- dependências do Maven",
        "|-- README.md",
        "`-- src/main/",
        "    |-- java/com/altzen/família/api/",
        "    |   |-- SistemaFamiliaApiApplication.java   <- main() do Spring Boot",
        "    |   |",
        "    |   |-- model/          <- Pessoa, Animal, Casa, Carro",
        "    |   |-- repository/     <- listas em memória (singleton)",
        "    |   |-- service/        <- regras de negócio (Services)",
        "    |   |-- controller/     <- endpoints REST (@RestController)",
        "    |   |-- dto/            <- Request / Response de cada rota",
        "    |   `-- exception/      <- BusinessException, NotFoundException,",
        "    |                          ErroResponse, GlobalExceptionHandler",
        "    `-- resources/",
        "        `-- application.properties   <- porta, prefixo /api, etc.",
    ]

    add_rect(slide, Inches(0.6), Inches(2.05), Inches(12.15), Inches(4.75), fill=SLATE_900, rounded=True, corner_radius=0.04)
    tb = slide.shapes.add_textbox(Inches(0.85), Inches(2.25), Inches(11.7), Inches(4.4))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.22
        r = p.add_run()
        r.text = line
        r.font.name = FONT_CODE
        r.font.size = Pt(13)
        if "<-" in line:
            r.font.color.rgb = SLATE_100
        else:
            r.font.color.rgb = SKY_400 if line.strip().endswith("/") else SLATE_100

    add_footer(slide, 21, TOTAL_SLIDES)


def slide_menu_endpoints():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Rotas da API", "Do menu antigo para os endpoints REST")

    header = ["#", "Opção do menu antigo", "Método", "Endpoint"]
    rows = [
        ["1", "Cadastrar pessoa", "POST", "/pessoas"],
        ["5", "Listar família", "GET", "/pessoas"],
        ["9", "Ver patrimônio", "GET", "/pessoas/{id}/patrimônio"],
        ["3", "Cadastrar casa", "POST", "/casas"],
        ["10", "Atribuir dono (casa)", "PUT", "/casas/{id}/dono"],
        ["11", "Transferir casa", "POST", "/casas/{id}/transferência"],
        ["4", "Cadastrar carro", "POST", "/carros"],
        ["12", "Transferir carro", "POST", "/carros/{id}/transferência"],
        ["14", "Pagar dívida", "POST", "/pessoas/{id}/pagamento-dívida"],
    ]

    left = Inches(0.6)
    top = Inches(2.05)
    total_w = Inches(12.15)
    col_widths = [Inches(0.7), Inches(4.5), Inches(1.5), Inches(5.45)]
    row_h = Inches(0.42)

    # header
    x = left
    for i, cell in enumerate(header):
        add_rect(slide, x, top, col_widths[i], row_h, fill=SLATE_800, rounded=False)
        align = PP_ALIGN.CENTER if i in (0, 2) else PP_ALIGN.LEFT
        pad = Inches(0.15) if align == PP_ALIGN.LEFT else Inches(0)
        add_text(slide, x + pad, top, col_widths[i] - pad, row_h,
                 cell, font=FONT_TITLE, size=12, bold=True, color=WHITE,
                 align=align, anchor=MSO_ANCHOR.MIDDLE)
        x = x + col_widths[i]

    # rows
    for ri, row in enumerate(rows):
        yy = top + row_h * (ri + 1)
        bg = SLATE_50 if ri % 2 == 0 else WHITE
        x = left
        for i, cell in enumerate(row):
            add_rect(slide, x, yy, col_widths[i], row_h, fill=bg, rounded=False)
            align = PP_ALIGN.CENTER if i in (0, 2) else PP_ALIGN.LEFT
            pad = Inches(0.15) if align == PP_ALIGN.LEFT else Inches(0)
            color = SLATE_700
            bold = False
            font = FONT_BODY
            if i == 2:  # metodo HTTP
                if cell == "GET": color = EMERALD_600
                elif cell == "POST": color = SKY_500
                elif cell == "PUT": color = AMBER_600
                elif cell == "DELETE": color = RED_600
                bold = True
                font = FONT_TITLE
            if i == 3:
                font = FONT_CODE
                color = SLATE_800
            add_text(slide, x + pad, yy, col_widths[i] - pad, row_h,
                     cell, font=font, size=11, bold=bold, color=color,
                     align=align, anchor=MSO_ANCHOR.MIDDLE)
            x = x + col_widths[i]

    add_footer(slide, 22, TOTAL_SLIDES)


def slide_testar_header():
    slide_section(
        "PARTE 6",
        "Testando a API",
        "Swagger, cURL e a padronização das respostas de erro.",
    )


def slide_swagger():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Testando", "Swagger UI - explore a API sem instalar nada")

    add_bullet_list(
        slide,
        Inches(0.6),
        Inches(2.1),
        Inches(6.0),
        Inches(4.5),
        [
            ("Rota: ", "http://localhost:8080/api/swagger-ui.html"),
            ("Lista todos os endpoints ", "agrupados por Controller."),
            ("Mostra os schemas ", "de cada Request e Response."),
            ("Permite executar ", "a requisição no navegador (sem Postman)."),
            ("Exemplos automáticos ", "do JSON esperado em cada rota."),
            ("Gerado pelo springdoc-openapi ", "a partir das anotações do código."),
        ],
        size=13,
    )

    # visual: janela de navegador
    x = Inches(6.9); y = Inches(2.1); w = Inches(5.9); h = Inches(4.35)
    add_rect(slide, x, y, w, h, fill=WHITE, line=SLATE_200, line_width=1.0, rounded=True, corner_radius=0.03)
    # barra do topo
    add_rect(slide, x, y, w, Inches(0.45), fill=SLATE_100, rounded=False)
    for i, c in enumerate([RGBColor(0xEF,0x44,0x44), RGBColor(0xF5,0x9E,0x0B), RGBColor(0x22,0xC5,0x5E)]):
        cx = x + Inches(0.2) + Inches(0.25) * i
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, y + Inches(0.13), Inches(0.20), Inches(0.20))
        dot.fill.solid(); dot.fill.fore_color.rgb = c; dot.line.fill.background()
    add_text(slide, x + Inches(1.2), y, w - Inches(1.4), Inches(0.45),
             "localhost:8080/api/swagger-ui.html",
             font=FONT_CODE, size=11, color=SLATE_600, anchor=MSO_ANCHOR.MIDDLE)

    # conteudo simulado do swagger
    y2 = y + Inches(0.65)
    endpoints = [
        ("GET", "/pessoas", EMERALD_600),
        ("POST", "/pessoas", SKY_500),
        ("GET", "/pessoas/{id}/patrimônio", EMERALD_600),
        ("POST", "/casas/{id}/transferência", SKY_500),
        ("PUT", "/animais/{id}/dono", AMBER_600),
    ]
    for i, (m, path, color) in enumerate(endpoints):
        yy = y2 + Inches(0.55) * i
        add_rect(slide, x + Inches(0.2), yy, Inches(0.75), Inches(0.42), fill=color, rounded=True, corner_radius=0.15)
        add_text(slide, x + Inches(0.2), yy, Inches(0.75), Inches(0.42),
                 m, font=FONT_TITLE, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(1.05), yy + Inches(0.02), Inches(4.7), Inches(0.4),
                 path, font=FONT_CODE, size=13, color=SLATE_700, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, 24, TOTAL_SLIDES)


def slide_curl():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Testando", "Chamada real com cURL: cadastrar pessoa")

    # Bloco 1: requisicao
    add_text(slide, Inches(0.6), Inches(2.05), Inches(6.0), Inches(0.4),
             "1) Enviamos um POST para /api/pessoas com um JSON:",
             font=FONT_BODY, size=13, bold=True, color=SLATE_700)
    add_code_block(
        slide, Inches(0.6), Inches(2.50), Inches(6.0), Inches(2.2),
        [
            "curl -X POST http://localhost:8080/api/pessoas \\",
            "  -H \"Content-Type: application/json\" \\",
            "  -d '{",
            "        \"nome\": \"Pai\",",
            "        \"idade\": 45,",
            "        \"parentesco\": \"Pai\",",
            "        \"carteira\": 5000",
            "      }'",
        ],
        size=11,
    )

    # Bloco 2: resposta
    add_text(slide, Inches(6.85), Inches(2.05), Inches(6.0), Inches(0.4),
             "2) A API responde 201 Created com o corpo:",
             font=FONT_BODY, size=13, bold=True, color=SLATE_700)
    add_code_block(
        slide, Inches(6.85), Inches(2.50), Inches(6.0), Inches(2.2),
        [
            "{",
            "  \"id\": 1,",
            "  \"nome\": \"Pai\",",
            "  \"idade\": 45,",
            "  \"parentesco\": \"Pai\",",
            "  \"carteira\": 5000.0,",
            "  \"dívida\": 0.0",
            "}",
        ],
        size=11,
    )

    # comentario final
    add_rect(slide, Inches(0.6), Inches(4.95), Inches(12.15), Inches(1.85), fill=SKY_50, rounded=True, corner_radius=0.04, line=SKY_100, line_width=0.75)
    add_text(slide, Inches(0.9), Inches(5.10), Inches(11.5), Inches(0.4),
             "O que aconteceu por baixo dos panos?",
             font=FONT_TITLE, size=15, bold=True, color=SKY_600)
    add_bullet_list(
        slide, Inches(0.9), Inches(5.55), Inches(11.5), Inches(1.3),
        [
            ("Jackson ", "transformou o JSON em PessoaRequest."),
            ("@Valid ", "checou @NotBlank, @Min, @PositiveOrZero."),
            ("Controller ", "chamou service.cadastrar(...), que salvou no Repository."),
            ("PessoaResponse.from(...) ", "converteu o Model em JSON e devolveu 201."),
        ],
        size=11, bullet_color=SKY_500,
    )

    add_footer(slide, 25, TOTAL_SLIDES)


def slide_erros():
    slide = add_slide()
    set_background(slide, WHITE)
    add_page_header(slide, "Testando", "Padronização de erros")

    add_text(
        slide, Inches(0.6), Inches(1.95), Inches(12.0), Inches(0.5),
        "Todo erro devolve um JSON no MESMO formato. O aluno não precisa adivinhar o que veio.",
        font=FONT_BODY, size=13, color=SLATE_600,
    )

    # tabela de status
    header = ["Código", "Quando acontece", "Vem de"]
    rows = [
        ["400", "Campo inválido no payload (@NotBlank, @Positive...)", "MethodArgumentNotValidException"],
        ["404", "Id de pessoa/animal/casa/carro não existe", "NotFoundException"],
        ["422", "Regra de negócio (menor de idade, sem saldo, etc.)", "BusinessException"],
        ["500", "Falha inesperada (bug, NullPointer...)", "qualquer outra Exception"],
    ]
    left = Inches(0.6); top = Inches(2.55)
    col_widths = [Inches(1.3), Inches(6.0), Inches(4.85)]
    row_h = Inches(0.42)
    colors = {"400": AMBER_600, "404": AMBER_600, "422": RED_600, "500": RED_600}

    x = left
    for i, cell in enumerate(header):
        add_rect(slide, x, top, col_widths[i], row_h, fill=SLATE_800, rounded=False)
        add_text(slide, x + Inches(0.15), top, col_widths[i] - Inches(0.15), row_h,
                 cell, font=FONT_TITLE, size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        x = x + col_widths[i]

    for ri, row in enumerate(rows):
        yy = top + row_h * (ri + 1)
        bg = SLATE_50 if ri % 2 == 0 else WHITE
        x = left
        for i, cell in enumerate(row):
            add_rect(slide, x, yy, col_widths[i], row_h, fill=bg, rounded=False)
            if i == 0:
                add_text(slide, x, yy, col_widths[i], row_h,
                         cell, font=FONT_TITLE, size=14, bold=True,
                         color=colors.get(cell, SLATE_700),
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            elif i == 2:
                add_text(slide, x + Inches(0.15), yy, col_widths[i] - Inches(0.15), row_h,
                         cell, font=FONT_CODE, size=11, color=SLATE_700, anchor=MSO_ANCHOR.MIDDLE)
            else:
                add_text(slide, x + Inches(0.15), yy, col_widths[i] - Inches(0.15), row_h,
                         cell, font=FONT_BODY, size=12, color=SLATE_700, anchor=MSO_ANCHOR.MIDDLE)
            x = x + col_widths[i]

    # exemplo real
    add_text(slide, Inches(0.6), Inches(4.75), Inches(12.0), Inches(0.35),
             "Exemplo de erro 422 quando falta saldo na transferência:",
             font=FONT_BODY, size=12, bold=True, color=SLATE_700)
    add_code_block(
        slide, Inches(0.6), Inches(5.15), Inches(12.15), Inches(1.85),
        [
            "{",
            "  \"timestamp\": \"2026-06-29T20:15:33.123\",",
            "  \"status\": 422,",
            "  \"erro\": \"Regra de negócio violada\",",
            "  \"mensagem\": \"Saldo insuficiente: faltam R$ 199000,00. Reenvie com 'aceitaEmprestimo=true'.\"",
            "}",
        ],
        size=10,
    )

    add_footer(slide, 26, TOTAL_SLIDES)


def slide_perguntas():
    slide = add_slide()
    set_background(slide, WHITE)

    # decoracao
    add_rect(slide, Inches(9.5), Inches(-1.0), Inches(6.0), Inches(4.0), fill=SKY_50, rounded=True, corner_radius=0.1)
    add_rect(slide, Inches(10.5), Inches(3.2), Inches(3.5), Inches(3.5), fill=SKY_100, rounded=True, corner_radius=0.1)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.4), Inches(0.6), Inches(1.6), Inches(1.6))
    circ.fill.solid(); circ.fill.fore_color.rgb = SKY_500; circ.line.fill.background()

    add_rich_text(
        slide,
        Inches(0.7), Inches(2.4), Inches(11.0), Inches(2.0),
        [[
            {"text": "Perguntas?", "font": FONT_TITLE, "size": 84, "bold": True, "color": SLATE_900},
        ]],
    )
    add_text(
        slide, Inches(0.7), Inches(4.55), Inches(11.0), Inches(0.6),
        "Obrigado pela atenção!",
        font=FONT_BODY, size=22, color=SLATE_500,
    )
    add_line(slide, Inches(0.7), Inches(5.35), Inches(4.5), Inches(5.35), color=SKY_500, width=2.0)
    add_text(
        slide, Inches(0.7), Inches(5.50), Inches(11.0), Inches(0.5),
        "Agora é sua vez de abrir o código e explorar cada camada.",
        font=FONT_BODY, size=15, bold=True, color=SKY_600,
    )
    add_text(
        slide, Inches(0.7), Inches(5.90), Inches(11.0), Inches(0.5),
        "Comece pelo endpoint POST /pessoas e siga o caminho: Controller -> Service -> Repository.",
        font=FONT_BODY, size=13, color=SLATE_500,
    )


# ----------------------------------------------------------------------------
# MONTAGEM DA APRESENTACAO
# ----------------------------------------------------------------------------

# 1
slide_capa()
# 2
slide_section(
    "PARTE 1",
    "O Ponto de Partida",
    "Antes de mergulhar no código, entendendo de onde saímos e para onde vamos.",
)
# 3
slide_de_onde_viemos()
# 4
slide_o_que_e_mvc()
# 5
slide_arquitetura_camadas()
# 6
slide_camadas_pratica_header()
# 7
slide_model()
# 8
slide_repository()
# 9
slide_service()
# 10
slide_controller()
# 11
slide_dto()
# 12
slide_exception()
# 13
slide_di_header()
# 14
slide_di_o_que_e()
# 15
slide_di_como_spring_faz()
# 16
slide_di_no_projeto()
# 17
slide_ciclo_header()
# 18
slide_ciclo_passo_a_passo()
# 19
slide_validacao()
# 20
slide_estrutura_header()
# 21
slide_estrutura_pastas()
# 22
slide_menu_endpoints()
# 23
slide_testar_header()
# 24
slide_swagger()
# 25
slide_curl()
# 26
slide_erros()
# 27
slide_perguntas()


import os
BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, "Descomplicando a Arquitetura - Sistema Familia API.pptx")
prs.save(OUT)
print(f"OK: gerado {OUT}")
print(f"Total de slides: {len(prs.slides)}")
