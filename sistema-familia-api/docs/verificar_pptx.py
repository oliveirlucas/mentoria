"""
Verifica o pptx gerado: quantos slides, quantos shapes por slide,
e detecta shapes que estao fora dos limites do slide.
"""
import os
from pptx import Presentation
from pptx.util import Emu

BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, "Descomplicando a Arquitetura - Sistema Familia API.pptx")
prs = Presentation(OUT)

SW = prs.slide_width
SH = prs.slide_height
print(f"Slide: {SW/914400:.2f} x {SH/914400:.2f} in")
print(f"Total de slides: {len(prs.slides)}")

issues = []
for i, slide in enumerate(prs.slides, 1):
    n_shapes = len(slide.shapes)
    for s in slide.shapes:
        try:
            l, t, w, h = s.left, s.top, s.width, s.height
        except Exception:
            continue
        if l is None or t is None or w is None or h is None:
            continue
        right = l + w
        bottom = t + h
        # tolera pequenos overflows (formas decorativas)
        if right > SW + Emu(10000) or bottom > SH + Emu(10000) or l < -Emu(200000) or t < -Emu(200000):
            # ignora as formas decorativas sabidas (bloco azul superior direito)
            issues.append((i, s.name, f"L={l/914400:.2f} T={t/914400:.2f} R={right/914400:.2f} B={bottom/914400:.2f}"))

    print(f"Slide {i:2d}: {n_shapes:2d} shapes")

print()
if issues:
    print("SHAPES POTENCIALMENTE FORA DO SLIDE (verificar visualmente):")
    for i, n, coord in issues:
        print(f"  Slide {i}: {n} -> {coord}")
else:
    print("Todas os shapes cabem no slide.")
